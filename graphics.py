import tkinter as tk
import customtkinter
from pptx import Presentation
from bardapi import BardCookies
import os

cookie_dict = {
    "__Secure-1PSID": "awi4iwZztiuw5Ce7U_0HS-MQ8dtI7zMInDnnIeDY8kpLcK7FZhKm4CqlujwZuRo5Ol9zbw.",
    "__Secure-1PSIDTS": "sidts-CjEBSAxbGWDY_K0OSUUPfYpiMtSsD9PRSQHpYueLOU9IBEQxf8D7lPVzoe-x-2r7T8IpEAA",
    "__Secure-1PSIDCC": "APoG2W_akEsNA8Xjzu-icq_aDg7a3R_qCv_YIqIPP6qm1z_HCaAELB6WJVQeoB1FO8AxGIVN4g"
}

bard = BardCookies(cookie_dict=cookie_dict)

customtkinter.set_appearance_mode("System")  # Modes: "System" (standard), "Dark", "Light")
customtkinter.set_default_color_theme("dark-blue")  # Themes: "blue" (standard), "green", "dark-blue"

os.environ['HTTP_PROXY'] = 'http://10.11.4.1:3128'
os.environ['HTTPS_PROXY'] = 'http://10.11.4.1:3128'


class App(customtkinter.CTk):
    def __init__(self):
        super().__init__()

        # configure window
        self.title("Powerfy")
        self.geometry(f"{460}x{520}")

        # configure grid layout (4x4)
        self.grid_columnconfigure(1, weight=1)
        self.grid_columnconfigure((2, 2), weight=1)
        self.grid_rowconfigure((0, 1, 2), weight=1)

        # create sidebar frame with widgets
        self.sidebar_frame = customtkinter.CTkFrame(self, corner_radius=0)
        self.sidebar_frame.grid(row=0, column=0, rowspan=5, sticky="nsew")
        self.sidebar_frame.grid_rowconfigure(5, weight=1)

        self.logo_label = customtkinter.CTkLabel(self.sidebar_frame, text="The Powerfy", font=customtkinter.CTkFont(size=20, weight="bold"))
        self.logo_label.grid(row=0, column=0, padx=20, pady=(20, 10))

        self.sidebar_button_1 = customtkinter.CTkButton(self.sidebar_frame, command=self.sidebar_button_event)
        self.sidebar_button_1.grid(row=1, column=0, padx=20, pady=10)

        self.sidebar_button_2 = customtkinter.CTkButton(self.sidebar_frame, command=self.sidebar_button_event)
        self.sidebar_button_2.grid(row=2, column=0, padx=20, pady=10)

        self.sidebar_button_3 = customtkinter.CTkButton(self.sidebar_frame, command=self.sidebar_button_event)
        self.sidebar_button_3.grid(row=3, column=0, padx=20, pady=10)

        # Theme settings
        self.appearance_mode_label = customtkinter.CTkLabel(self.sidebar_frame, text="Appearance Mode", anchor="w")
        self.appearance_mode_label.grid(row=5, column=0, padx=20, pady=(1, 0))
        self.appearance_mode_optionemenu = customtkinter.CTkOptionMenu(self.sidebar_frame, values=["Light", "Dark", "System"],command=self.change_appearance_mode_event)
        self.appearance_mode_optionemenu.grid(row=6, column=0, padx=20, pady=(5, 10))

        # UI settings
        self.scaling_label = customtkinter.CTkLabel(self.sidebar_frame, text="UI Scaling", anchor="w")
        self.scaling_label.grid(row=7, column=0, padx=20, pady=(10, 0))
        self.scaling_optionemenu = customtkinter.CTkOptionMenu(self.sidebar_frame, values=["80%", "90%", "100%", "110%", "130%", "150%"],command=self.change_scaling_event)
        self.scaling_optionemenu.grid(row=8, column=0, padx=20, pady=(1, 20))

        # Tabview Buttons
        self.tabview = customtkinter.CTkTabview(self, corner_radius=5)
        self.tabview.grid(row=0, column=1, padx=20, pady=20, sticky="nsew")
        self.tabview.add("Settings")
        self.tabview.add("Options")
        self.tabview.add("General")

        self.tabview.tab("Settings").grid_columnconfigure(0, weight=1)
        self.tabview.tab("Options").grid_columnconfigure(0, weight=1)
        self.tabview.tab("General").grid_columnconfigure(0, weight=1)

        # Settings Tab
        self.combobox_1 = customtkinter.CTkComboBox(self.tabview.tab("Settings"), values=["5", "7", "10"])
        self.combobox_1.grid(row=0, column=0, padx=20, pady=(20, 10))

        self.combobox_2 = customtkinter.CTkComboBox(self.tabview.tab("Settings"), values=["English", "French", "Romanian"])
        self.combobox_2.grid(row=1, column=0, padx=20, pady=(15, 10))

        self.combobox_3 = customtkinter.CTkComboBox(self.tabview.tab("Settings"), values=["History", "Literature", "Technology", "Biology"])
        self.combobox_3.grid(row=2, column=0, padx=20, pady=(15, 10))

        # Create a frame for entry points and button
        entry_frame = customtkinter.CTkFrame(self,corner_radius=5)
        entry_frame.grid(row=1, column=1, rowspan=4, padx=20, pady=(10, 20), sticky="nsew")

        # Entry Points
        self.entry1_label = customtkinter.CTkLabel(entry_frame, text="Project Name", anchor="w", font=("Arial", 12))
        self.entry1_label.grid(row=0, column=0, padx=10, pady=(3,0), sticky="nsew")

        self.entry1 = customtkinter.CTkEntry(entry_frame, placeholder_text="")
        self.entry1.grid(row=1, column=0, columnspan=1, padx=20, pady=5, sticky="nsew")

        self.entry2_label = customtkinter.CTkLabel(entry_frame, text="Optional", anchor="w", font=("Arial", 12))
        self.entry2_label.grid(row=2, column=0, padx=10, pady=(3,0), sticky="nsew")

        self.entry2 = customtkinter.CTkEntry(entry_frame, placeholder_text="")
        self.entry2.grid(row=3, column=0, columnspan=1, padx=20, pady=5, sticky="nsew")

        # Button
        self.main_button_1 = customtkinter.CTkButton(text='Generate',master=self, command=self.open_input_dialog_event)
        self.main_button_1.grid(row=5, column=1, padx=50, pady=(5, 20), sticky="nsew")

        # Default Values
        self.sidebar_button_3.configure(state="disabled", text="Development")
        self.sidebar_button_2.configure(state="disabled", text="Social")
        self.sidebar_button_1.configure(state="disabled", text="Operations")
        self.appearance_mode_optionemenu.set("Dark")
        self.scaling_optionemenu.set("100%")
        self.combobox_1.set("Slides")
        self.combobox_2.set("Language")
        self.combobox_3.set("Templates")

    def open_input_dialog_event(self):

        project_name = self.entry1.get()
        optional_text = self.entry2.get()

        combobox_value_1 = self.combobox_1.get()
        combobox_value_2 = self.combobox_2.get()
        combobox_value_3 = self.combobox_3.get()

        print("Project Name:", project_name)
        print("Optional:", optional_text)
        print("Combobox 1:", combobox_value_1)
        print("Combobox 2:", combobox_value_2)
        print("Combobox 3:", combobox_value_3)

        path = 'templates/' + combobox_value_3 +".pptx"
        presentation_path = os.path.join(path)
        presentation = Presentation(presentation_path)

        new_title = combobox_value_2
        new_subtitle = combobox_value_3

        slide0 = presentation.slides[0]

        if slide0.shapes.title:
            slide0.shapes.title.text = new_title
        if len(slide0.placeholders) > 1:
            slide0.placeholders[1].text = new_subtitle

        output_path = os.path.join(path)
        presentation.save(output_path)

        #----------------------------------------------------------------

        if combobox_value_1 == 5:
            question = f"Can write an presentation with 5 slides about {project_name} in the following order: "

        #----------------------------------------------------------------

        answer = bard.get_answer(str(project_name))['content']

        #----------------------------------------------------------------

        print(f"Modified presentation saved to: {output_path}")
        print(answer)

    def change_appearance_mode_event(self, new_appearance_mode: str):
        customtkinter.set_appearance_mode(new_appearance_mode)
    def change_scaling_event(self, new_scaling: str):
        new_scaling_float = int(new_scaling.replace("%", "")) / 100
        customtkinter.set_widget_scaling(new_scaling_float)
    def sidebar_button_event(self):
        print("sidebar_button click")

if __name__ == "__main__":
    app = App()
    app.mainloop()
